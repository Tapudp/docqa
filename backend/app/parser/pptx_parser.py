import io

from app.parser.base import DocumentParser, ParseResult

_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


class PptxParser(DocumentParser):
    def supports(self, mime_type: str) -> bool:
        return mime_type == _PPTX_MIME

    async def parse(self, data: bytes, mime_type: str) -> ParseResult:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(data))
        pages: list[str] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            lines: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            lines.append(text)
            pages.append(f"[Slide {slide_num}]\n" + "\n".join(lines) if lines else f"[Slide {slide_num}]")

        if not pages:
            pages = [""]

        return ParseResult(total_pages=len(pages), page_texts=pages, failed_pages=[])
